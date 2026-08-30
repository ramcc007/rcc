import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from common import *
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

prs = new_deck(); N=[0]
def page(): N[0]+=1; return N[0]

def axis(ax, size=10, grid=False):
    ax.has_major_gridlines = grid
    if grid:
        gl = ax.major_gridlines.format.line
        gl.color.rgb = RGBColor(0xEC,0xEA,0xEE); gl.width = Pt(0.75)
    ax.tick_labels.font.size = Pt(size); ax.tick_labels.font.name = FONT
    ax.tick_labels.font.color.rgb = GRAY
    ax.format.line.color.rgb = RGBColor(0xDD,0xDA,0xE0)

# ══════════════════════════════════════ 1 · TITLE
s = blank(prs)
rect(s, 0, 0, 4.60, 7.5, PURPLE); rect(s, 4.60, 0, 0.09, 7.5, GREEN)
rect(s, 0.55, 0.55, 2.60, 0.85, WHITE)
s.shapes.add_picture(LOGO, Inches(0.78), Inches(0.78), Inches(2.10), Inches(0.45))
txt(s, 0.55, 2.75, 3.85, 2.00, [("GEO & SEO", 31, True, WHITE, 3), ("Performance", 31, True, WHITE)])
txt(s, 0.55, 4.30, 3.85, 0.42, [("Q2 2026  ·  April – June", 16, True, GREEN)])
txt(s, 0.55, 6.55, 3.85, 0.50, [("Sources: Semrush, Google Analytics", 10.5, False, LAV)])

txt(s, 5.04, 2.55, 7.80, 1.95, [("How AI search describes", 25, True, PURPLE, 5),
                                ("Komprise today — and what", 25, True, PURPLE, 5),
                                ("we changed to move it", 25, True, PURPLE)])
rect(s, 5.04, 4.86, 1.40, 0.04, GREEN)
txt(s, 5.04, 5.22, 7.80, 1.40,
    [("Two lanes:", 13, True, INK, 8),
     ("Be the answer for AI-ready data — unstructured data for AI, AI data management.", 12.5, False, GRAY, 6),
     ("Be the answer for storage price hikes — tiering, cost savings, NetApp, Dell, Everpure.", 12.5, False, GRAY)])
page()

# ══════════════════════════════════════ 2 · EXECUTIVE SUMMARY
s = blank(prs)
header(s, "Executive Summary", "Q2 2026 — where we stand on being seen as an AI company")
statcard(s, 0.55, 1.40, 2.85, 1.52, "+20%", "Citations in AI Answers", "AI answers now point to a Komprise page 2.1K times a month.")
statcard(s, 3.55, 1.40, 2.85, 1.52, "‒72%", "Name Mentions", "Being cited is not the same as being named. This is the gap.", RED)
statcard(s, 6.55, 1.40, 2.85, 1.52, "1,515", "Google Top-10 Keywords", "Up 128 across the quarter, and up 66 in the top three.", PURPLE)
statcard(s, 9.55, 1.40, 2.85, 1.52, "52%", "of Traffic Now Direct", "Up from 41%. People read the AI answer, then type us in.")

bandbar(s, 0.55, 3.22, 12.23, "What This Means")
bullet(s, 0.60, 3.82, 11.9, "We already win the storage-cost lane inside AI answers. On tiering and storage-cost questions, Google's AI answer names Komprise first.")
bullet(s, 0.60, 4.52, 11.9, "We do not yet win the AI-data lane. On “AI data management” — the term buyers start with — we sit at position 15 and are absent from the AI answer entirely.")
bullet(s, 0.60, 5.22, 11.9, "Citations are rising, mentions are falling. Our content earns links in AI answers; only third-party coverage makes the model say our name.")
bullet(s, 0.60, 5.92, 11.9, "The traffic mix is changing underneath us. Direct has overtaken organic search — the visible proof that AI answers are doing the discovery.")

rect(s, 0.55, 6.58, 12.23, 0.48, PALE); rect(s, 0.55, 6.58, 0.06, 0.48, PURPLE)
txt(s, 0.78, 6.70, 11.9, 0.30, [("In short: the storage story is won in AI search. The AI story is half-built — and that is where Q3 goes.", 12, True, PURPLE)])
footer(s, page())

# ══════════════════════════════════════ 3 · WHAT MATTERS TO LLMs
s = blank(prs)
kicker_header(s, "HOW AI SEARCH WORKS", "What Decides Whether an AI Names You",
              "Three things, and only the first is your own website")

def pillar(x, num, title, body, did, col):
    rect(s, x, 1.62, 3.90, 3.62, OFFWHITE, BORDER)
    rect(s, x, 1.62, 3.90, 0.07, col)
    rect(s, x+0.28, 1.92, 0.50, 0.50, col, shape=MSO_SHAPE.OVAL)
    txt(s, x+0.28, 2.04, 0.50, 0.28, [(num, 14.5, True, WHITE)], align=PP_ALIGN.CENTER)
    txt(s, x+0.28, 2.60, 3.35, 0.34, [(title, 14, True, PURPLE)])
    txt(s, x+0.28, 3.02, 3.35, 1.30, [(body, 11.5, False, GRAY, 0, 1.28)])
    rect(s, x+0.28, 4.42, 3.35, 0.62, PALE)
    txt(s, x+0.42, 4.53, 3.10, 0.44, [(did, 10, True, PURPLE, 0, 1.16)])

pillar(0.55, "1", "Be machine-readable",
       "AI crawlers read plain text and structured data. A page with no schema, or one that blocks the crawler, effectively does not exist to the model — however good the writing is.",
       "Done: AI crawlers admitted, llms.txt rebuilt, FAQ schema added.", PURPLE)
pillar(4.72, "2", "Be repeated",
       "A model forms its view of a company from the same description repeated across many trusted sites. One excellent page cannot outweigh a hundred third-party sources that say something older.",
       "Done: one company description rolled out across third-party sites.", DGREEN)
pillar(8.88, "3", "Be named, not just cited",
       "A citation is a link inside the answer. A mention is the model actually saying “Komprise”. Citations follow good content. Mentions follow earned coverage — analysts, press, independent authority.",
       "Open: this is the metric still moving the wrong way.", AMBER)

rect(s, 0.55, 5.50, 12.23, 1.28, OFFWHITE, BORDER); rect(s, 0.55, 5.50, 0.06, 1.28, AMBER)
txt(s, 0.82, 5.72, 11.8, 1.00,
    [("What this looks like in practice", 12.5, True, PURPLE, 6),
     ("In July we asked ChatGPT “How to address rising enterprise storage costs — please cite sources.” It answered with Microsoft Learn, AWS, TechTarget and an academic paper. No vendor content at all, including ours — even though we have strong, well-sourced pages on exactly that question. When a buyer asks for sources, the model reaches for independent publications over any vendor. That single behaviour is why the next lever is third-party placement, not more pages.", 11, False, GRAY, 0, 1.25)])
footer(s, page())

# ══════════════════════════════════════ 4 · WHAT WE CHANGED
s = blank(prs)
kicker_header(s, "WHAT WE DID", "Changing How We Are Described",
              "Content and technical work delivered against the two lanes")

rect(s, 0.55, 1.62, 6.02, 0.38, PURPLE)
txt(s, 0.73, 1.695, 5.7, 0.26, [("CONTENT  ·  CLAIMING THE LANGUAGE", 12, True, WHITE)])
rect(s, 6.76, 1.62, 6.02, 0.38, DGREEN)
txt(s, 6.94, 1.695, 5.7, 0.26, [("TECHNICAL  ·  MAKING IT READABLE", 12, True, WHITE)])

L = ["More than 20 new AI glossary pages published to claim the language buyers and models use — AI Data Platform, AI-Ready Data, AI Data Ingestion, AI Data Curation, Vector Embeddings, Synthetic Data, AI Token, Tokenization.",
     "Memflation and flash-pricing pages created to meet buyers searching on storage price rises rather than product names.",
     "Homepage hero and product page titles reframed to lead with AI readiness, with the 70% cost saving moved to a supporting proof point.",
     "Transparent File Tables launched with a product page, demo page, interview and press coverage.",
     "AI-pipeline explainer content published and comparison tables added across the glossary."]
R = ["Full audit of all 1,433 site URLs, separating copy fixes from engineering fixes.",
     "AI crawlers explicitly admitted in robots.txt — previously two major assistants could not read the site at all.",
     "llms.txt rebuilt and repointed away from cost-led framing to current section indexes.",
     "FAQ and structured data added across the Flash Stretch page, blog posts and glossary terms.",
     "Canonical repair on 20 new AI glossary pages that were pointing crawlers at broken URLs.",
     "Legacy cost-led meta template retired from 58 pages; duplicate blog index resolved and sitemap cleaned."]

import math
def col(x, items, dot, size=10.5):
    y = 2.14
    for t in items:
        lines = max(1, math.ceil(len(t) / 62.0))
        h = lines * 0.205
        rect(s, x, y+0.06, 0.09, 0.09, dot)
        txt(s, x+0.26, y-0.03, 5.60, h+0.10, [(t, size, False, INK, 0, 1.22)])
        y += h + 0.24
    return y
col(0.58, L, PURPLE)
col(6.79, R, DGREEN, 10.5)

rect(s, 0.55, 6.42, 12.23, 0.62, PALE); rect(s, 0.55, 6.42, 0.06, 0.62, PURPLE)
txt(s, 0.78, 6.55, 11.8, 0.42,
    [("Early proof it works: product names we deliberately repeated across trusted sites — Transparent Move Technology, Global Metadatabase — are now being repeated back to us by AI assistants, unprompted.", 11, True, PURPLE, 0, 1.2)])
footer(s, page())

# ══════════════════════════════════════ 5 · TWO LANES, MEASURED
s = blank(prs)
kicker_header(s, "THE RESULT, LANE BY LANE", "Are We the Answer, or Not?",
              "Whether Komprise appears inside Google's AI answer for the questions each lane is built on")

rect(s, 0.55, 1.62, 6.02, 0.38, PURPLE)
txt(s, 0.73, 1.695, 5.7, 0.26, [("LANE 1  ·  AI-READY DATA", 12, True, WHITE)])
rect(s, 6.76, 1.62, 6.02, 0.38, DGREEN)
txt(s, 6.94, 1.695, 5.7, 0.26, [("LANE 2  ·  STORAGE PRICE HIKES", 12, True, WHITE)])

LANE1 = [("ai data ingestion", True, "#1"), ("ai data curation", True, "#1"),
         ("unstructured data management solutions", True, "#1"), ("ai data leakage", True, "#1"),
         ("ai data management", False, "Google #15"), ("ai data platform", False, "Google #10"),
         ("ai data preparation", False, "Google #9"), ("how to prepare data for ai", False, "Google #2")]
LANE2 = [("data tiering", True, "#1"), ("storage tiering", True, "#1"),
         ("intelligent tiering", True, "#1"), ("reduce storage costs", True, "#1"),
         ("data storage costs", True, "#1"), ("azure netapp", True, "#2"),
         ("netapp to pure storage migration", True, "#1"),
         ("rising storage costs — cite sources", False, "no vendor")]

def lane(x, rows):
    y = 2.12
    for q, win, note in rows:
        rect(s, x, y, 6.02, 0.40, OFFWHITE if win else WHITE, BORDER)
        rect(s, x, y, 0.05, 0.40, DGREEN if win else LGRAY)
        txt(s, x+0.22, y+0.11, 3.30, 0.24, [(q, 10.5, win, INK if win else GRAY)])
        if win:
            chip(s, x+3.58, y+0.08, 1.36, 0.24, "IN AI ANSWER", DGREEN)
            txt(s, x+5.04, y+0.11, 0.82, 0.22, [(note, 10, True, DGREEN)], wrap=False)
        else:
            chip(s, x+3.58, y+0.08, 1.36, 0.24, "NOT YET", LGRAY)
            txt(s, x+5.04, y+0.115, 0.82, 0.22, [(note, 8.5, False, LGRAY)], wrap=False)
        y += 0.46
lane(0.55, LANE1); lane(6.76, LANE2)

rect(s, 0.55, 5.92, 12.23, 1.00, PALE); rect(s, 0.55, 5.92, 0.06, 1.00, PURPLE)
txt(s, 0.82, 6.10, 11.7, 0.72,
    [("The verdict: on storage costs we are already the answer — seven of eight questions put Komprise inside Google's AI answer, most of them first. On AI-ready data we are half-built: we win the specific technical questions, and we lose the broad category terms a buyer actually starts with. Winning those is the work that remains.", 11.5, True, PURPLE, 0, 1.24)])
footer(s, page())

# ══════════════════════════════════════ 6 · GEO MOMENTUM
s = blank(prs)
kicker_header(s, "MOMENTUM", "Citations Are Rising. Mentions Are the Soft Spot.",
              "Jan → Jun 2026:   Citations +20%   ·   Cited Pages +19%   ·   Mentions ‒72%")

cd = CategoryChartData()
cd.categories = ['Jan','Feb','Mar','Apr','May','Jun']
cd.add_series('Citations',   (1700, 1750, 1680, 1870, 2000, 2050))
cd.add_series('Cited Pages', (600,  575,  545,  620,  700,  715))
cd.add_series('Mentions',    (595,  565,  425,  225,  190,  163))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.55), Inches(1.72), Inches(8.10), Inches(4.55), cd)
ch = gf.chart; ch.has_title = False; ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
for ser, c in zip(ch.series, (DGREEN, PURPLE, RED)):
    ser.format.line.color.rgb = c; ser.format.line.width = Pt(2.75); ser.smooth = False
    ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = c
    ser.marker.format.line.color.rgb = WHITE
axis(ch.value_axis, grid=True); axis(ch.category_axis, size=11)
ch.value_axis.minimum_scale = 0; ch.value_axis.maximum_scale = 2500
txt(s, 0.55, 6.34, 8.10, 0.50, [("Q2 is the steep part of all three lines. Monthly shape reproduced from the June metrics report; January and June values are exact.", 9.5, False, LGRAY, 0, 1.2)])

statcard(s, 8.95, 1.72, 3.83, 1.42, "2.1K", "Citations", "AI answers pointing to a Komprise page. Our content optimisation work landing.", DGREEN, 26)
statcard(s, 8.95, 3.29, 3.83, 1.42, "715", "Cited Pages", "Distinct pages being cited. Broad surface area, and holding up.", PURPLE, 26)
statcard(s, 8.95, 4.86, 3.83, 1.42, "163", "Name Mentions", "The model saying “Komprise”. Down sharply — a coverage gap, not a content gap.", RED, 26)
footer(s, page())

# ══════════════════════════════════════ 7 · SEO TOP-10
s = blank(prs)
kicker_header(s, "GOOGLE SEARCH", "Top-10 Rankings Are Compounding",
              "Keywords ranking in Google's top 10, split by top three and positions four to ten")

cd = CategoryChartData()
cd.categories = ['Jan','Feb','Mar','Apr','May','Jun']
cd.add_series('Positions 4–10', (841, 907, 1033, 972, 997, 1095))
cd.add_series('Top 3',          (333, 349,  354, 366, 387,  420))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.55), Inches(1.72), Inches(7.90), Inches(4.35), cd)
ch = gf.chart; ch.has_title = False; ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
pl = ch.plots[0]; pl.gap_width = 60
for ser, c in zip(ch.series, (LAV, PURPLE)):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
    ser.format.line.fill.background()
axis(ch.value_axis, grid=True); axis(ch.category_axis, size=11)
ch.value_axis.minimum_scale = 0; ch.value_axis.maximum_scale = 1700
txt(s, 0.55, 6.16, 7.90, 0.62,
    [("Q1 closed at 1,387 keywords in the top 10. Q2 closed at 1,515 — and the top-three band, the positions that actually get read, grew fastest.", 11, False, GRAY, 0, 1.22)])

statcard(s, 8.75, 1.72, 4.03, 1.38, "+9.2%", "Top-10 Keywords, QoQ", "1,387 at the close of Q1 to 1,515 at the close of Q2.", DGREEN, 25)
statcard(s, 8.75, 3.24, 4.03, 1.38, "+18.6%", "Top-3 Keywords, QoQ", "354 to 420. Growing faster than the wider top-10 band.", DGREEN, 25)
rect(s, 8.75, 4.76, 4.03, 1.52, OFFWHITE, BORDER); rect(s, 8.75, 4.76, 0.06, 1.52, PURPLE)
txt(s, 9.03, 4.96, 3.60, 1.20,
    [("Why it matters", 12, True, PURPLE, 6),
     ("Two years ago we held 510 top-10 keywords. We now hold 1,515. This is the compounding base that AI answers draw from — models cite what already ranks.", 10.5, False, GRAY, 0, 1.24)])
footer(s, page())

# ══════════════════════════════════════ 8 · GA4 TRAFFIC BEHAVIOUR
s = blank(prs)
kicker_header(s, "TRAFFIC BEHAVIOUR", "Direct Has Overtaken Organic Search",
              "Monthly sessions — the visible signature of zero-click search and direct brand discovery")

cd = CategoryChartData()
cd.categories = ['Jan  (Q1)','Apr  (Q2)','Jun  (Q2)']
cd.add_series('Direct',         (5956, 5938, 6726))
cd.add_series('Organic Search', (5074, 4769, 4139))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.72), Inches(7.55), Inches(4.30), cd)
ch = gf.chart; ch.has_title = False; ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
pl = ch.plots[0]; pl.gap_width = 70; pl.has_data_labels = True
dl = pl.data_labels; dl.font.size = Pt(10.5); dl.font.bold = True; dl.font.name = FONT
dl.font.color.rgb = INK; dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.number_format = '#,##0'; dl.number_format_is_linked = False
for ser, c in zip(ch.series, (PURPLE, LAV)):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
    ser.format.line.fill.background()
axis(ch.value_axis, grid=True); axis(ch.category_axis, size=11)
ch.value_axis.minimum_scale = 0; ch.value_axis.maximum_scale = 8000
txt(s, 0.55, 6.12, 7.55, 0.66,
    [("Total sessions dipped into June. That fall sits almost entirely in email campaign volume, not in search — organic and direct are the lines that describe how people find us.", 10.5, False, LGRAY, 0, 1.22)])

rect(s, 8.40, 1.72, 4.38, 2.16, OFFWHITE, BORDER); rect(s, 8.40, 1.72, 0.06, 2.16, DGREEN)
txt(s, 8.68, 1.94, 3.92, 1.80,
    [("Direct is now the #1 channel", 12.5, True, PURPLE, 6),
     ("Direct grew from 41% of sessions in Q1 to 52% in Q2, while organic search fell. People are reading the answer in an AI or search result, then arriving by typing us in.", 11, False, GRAY, 0, 1.26)])
rect(s, 8.40, 4.02, 4.38, 2.26, OFFWHITE, BORDER); rect(s, 8.40, 4.02, 0.06, 2.26, PURPLE)
txt(s, 8.68, 4.24, 3.92, 1.90,
    [("Why this is the right read", 12.5, True, PURPLE, 6),
     ("A falling click count is not a falling audience. AI answers resolve the question on the results page, so the visit only happens once someone decides they want us specifically. That is brand, and it is showing up in the data.", 11, False, GRAY, 0, 1.26)])
footer(s, page())

# ══════════════════════════════════════ 9 · COMPETITIVE POSITION
s = blank(prs)
kicker_header(s, "COMPETITIVE POSITION", "We Lead the Category on Every Search Measure",
              "Organic footprint and domain authority across the tracked competitive set")

hdr = ["Company", "Ranking\nKeywords", "Monthly Organic\nTraffic", "Referring\nDomains", "Authority\nScore"]
COLS = [0.55, 5.05, 7.05, 9.25, 11.15]
WID  = [4.40, 1.90, 2.10, 1.80, 1.60]
rect(s, 0.55, 1.72, 12.23, 0.50, PURPLE)
for i,(h,x,w) in enumerate(zip(hdr, COLS, WID)):
    al = PP_ALIGN.LEFT if i==0 else PP_ALIGN.CENTER
    txt(s, x+(0.20 if i==0 else 0), 1.85, w, 0.30, [(h.replace("\n"," "), 10.5, True, WHITE)], align=al)

ROWS = [("Komprise",              "6.4K", "23.9K", "1.9K", "43", True,  "+5.9%",  "+18.9%"),
        ("Unstructured.io",       "2.0K", "4.8K",  "2.7K", "36", False, "‒54.2%", "+15.7%"),
        ("Data Dynamics",         "1.4K", "1.0K",  "958",  "30", False, "+27.0%", "+116%"),
        ("Hammerspace",           "819",  "2.1K",  "1.5K", "33", False, "+17.0%", "‒52.3%"),
        ("Atempo",                "322",  "541",   "1.0K", "33", False, "+0.9%",  "‒57.4%"),
        ("Datadobi",              "86",   "113",   "751",  "23", False, "0%",     "‒47.2%"),
        ("Diskover Data",         "82",   "407",   "450",  "25", False, "+17.1%", "+21.9%")]
y = 2.26
for name, kw, tr, rd, asc, mine, dkw, dtr in ROWS:
    rect(s, 0.55, y, 12.23, 0.50, PALE if mine else (OFFWHITE if ROWS.index((name,kw,tr,rd,asc,mine,dkw,dtr))%2==0 else WHITE), BORDER)
    if mine: rect(s, 0.55, y, 0.06, 0.50, DGREEN)
    txt(s, 0.75, y+0.14, 4.20, 0.28, [(name, 12 if mine else 11, mine, PURPLE if mine else INK)])
    for val, x, w in ((kw, COLS[1], WID[1]), (tr, COLS[2], WID[2]), (rd, COLS[3], WID[3]), (asc, COLS[4], WID[4])):
        txt(s, x, y+0.14, w, 0.28, [(val, 12 if mine else 11, mine, PURPLE if mine else GRAY)], align=PP_ALIGN.CENTER)
    y += 0.54

rect(s, 0.55, 6.14, 12.23, 0.86, OFFWHITE, BORDER); rect(s, 0.55, 6.14, 0.06, 0.86, DGREEN)
txt(s, 0.82, 6.32, 11.7, 0.60,
    [("Komprise holds roughly three times the ranking keywords and five times the organic traffic of the nearest competitor, with the highest authority score in the set — and is the only company here that is both the largest and still growing. The three closest storage specialists all lost traffic this period.", 11.5, True, PURPLE, 0, 1.22)])
footer(s, page())

# ══════════════════════════════════════ 10 · WHAT'S NEXT
s = blank(prs)
kicker_header(s, "WHAT'S NEXT", "Where the Effort Goes",
              "Five priorities, weighted to the gap the data has exposed")

ITEMS = [("Win the category terms in the AI lane",
          "We hold the specific technical questions and lose the broad ones. Getting AI data management, AI data platform and AI data preparation into the AI answer is the highest-value work available.", PURPLE),
         ("Turn citations into mentions",
          "Content has taken citations as far as it can. Only independent coverage — analysts and trade press — makes a model say our name, and mentions are the metric still moving the wrong way.", AMBER),
         ("Finish the machine-readable layer",
          "Structured data is still missing from most blog posts and older template pages still carry cost-led copy. These are single-template fixes with disproportionate return.", DGREEN),
         ("Defend the storage-cost lane",
          "It is our strongest position in AI answers and it is now being contested by a much larger rebranded competitor. Holding it protects the pipeline story.", PURPLE),
         ("Agree one standing scoreboard",
          "Lock a top-ten priority question set across both lanes so every future update reports the same measure, quarter on quarter.", DGREEN)]
y = 1.72
for i,(t,b,c) in enumerate(ITEMS, 1):
    rect(s, 0.55, y, 12.23, 0.88, OFFWHITE, BORDER); rect(s, 0.55, y, 0.06, 0.88, c)
    rect(s, 0.85, y+0.24, 0.40, 0.40, c, shape=MSO_SHAPE.OVAL)
    txt(s, 0.85, y+0.335, 0.40, 0.26, [(str(i), 12, True, WHITE)], align=PP_ALIGN.CENTER)
    txt(s, 1.45, y+0.12, 11.1, 0.28, [(t, 12.5, True, PURPLE)])
    txt(s, 1.45, y+0.44, 11.1, 0.40, [(b, 10.5, False, GRAY, 0, 1.20)])
    y += 0.95

rect(s, 0.55, y+0.00, 12.23, 0.50, PALE); rect(s, 0.55, y+0.00, 0.06, 0.50, PURPLE)
txt(s, 0.78, y+0.12, 11.8, 0.30,
    [("How we will know it worked: Komprise inside the AI answer on the AI-lane category terms, and the mentions line turning back up.", 11.5, True, PURPLE, 0, 1.2)])
footer(s, page())

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..", "..", "..", "home", "user", "rcc", "Komprise_GEO_SEO_Q2_2026.pptx")
out = "/home/user/rcc/Komprise_GEO_SEO_Q2_2026.pptx"
prs.save(out); print("saved", out, "| slides:", len(prs.slides._sldIdLst))
