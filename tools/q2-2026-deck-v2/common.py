import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Brand tokens lifted from the June 2026 metrics deck ----
PURPLE   = RGBColor(0x59, 0x25, 0x7B)
GREEN    = RGBColor(0x9F, 0xC4, 0x06)
DGREEN   = RGBColor(0x7A, 0x99, 0x05)
INK      = RGBColor(0x26, 0x24, 0x28)
GRAY     = RGBColor(0x59, 0x5A, 0x5B)
LGRAY    = RGBColor(0x8A, 0x8A, 0x8D)
RED      = RGBColor(0xB0, 0x3A, 0x2E)
AMBER    = RGBColor(0xC8, 0x8A, 0x1E)
LAV      = RGBColor(0xD6, 0xC4, 0xE3)
PALE     = RGBColor(0xE9, 0xE2, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF7, 0xF5, 0xF9)
BORDER   = RGBColor(0xE3, 0xE0, 0xE6)
FONT = "Open Sans"

LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "komprise-logo.png")
FOOTER = "Komprise  —  GEO & SEO Performance, Q2 2026"

def new_deck():
    p = Presentation()
    p.slide_width  = Emu(12192000)
    p.slide_height = Emu(6858000)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, shape=MSO_SHAPE.RECTANGLE, adj=None):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if adj is not None:
        try: sh.adjustments[0] = adj
        except Exception: pass
    return sh

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs = list of (text, size, bold, color, space_after_pt) or (text,size,bold,color)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    first = True
    for r in runs:
        text, size, bold, color = r[0], r[1], r[2], r[3]
        sa = r[4] if len(r) > 4 else 0
        ls = r[5] if len(r) > 5 else 1.15
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(sa)
        p.line_spacing = ls
        run = p.add_run(); run.text = text
        f = run.font
        f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    return tb

def header(s, title, sub=None, kicker=None):
    rect(s, 0, 0, 13.333, 0.09, PURPLE)
    rect(s, 0, 0.09, 1.60, 0.04, GREEN)
    s.shapes.add_picture(LOGO, Inches(11.55), Inches(0.30), Inches(1.42), Inches(0.30))
    y = 0.30
    if kicker:
        txt(s, 0.55, 0.26, 10.0, 0.28, [(kicker, 10.5, True, PURPLE)])
        y = 0.56
        txt(s, 0.55, y, 10.6, 0.50, [(title, 24, True, INK)])
        y2 = y + 0.52
    else:
        txt(s, 0.55, y, 10.6, 0.60, [(title, 25, True, PURPLE)])
        y2 = 0.88
    if sub:
        txt(s, 0.55, y2, 11.2, 0.32, [(sub, 12.5, False, GRAY)])
    return y2 + 0.45

def footer(s, n):
    txt(s, 0.55, 7.14, 7.0, 0.25, [(FOOTER, 9, False, LGRAY)])
    txt(s, 11.90, 7.14, 0.90, 0.25, [(str(n), 9, False, LGRAY)], align=PP_ALIGN.RIGHT)

def statcard(s, x, y, w, h, big, label, note, accent=DGREEN, bigsize=30):
    rect(s, x, y, w, h, OFFWHITE, BORDER)
    rect(s, x, y, 0.06, h, accent)
    txt(s, x+0.28, y+0.18, w-0.48, 0.50, [(big, bigsize, True, accent)])
    txt(s, x+0.28, y+0.18+0.48, w-0.48, 0.28, [(label, 10.5, True, PURPLE)])
    txt(s, x+0.28, y+0.18+0.80, w-0.48, 0.60, [(note, 9, False, GRAY, 0, 1.2)])

def bandbar(s, x, y, w, label):
    rect(s, x, y, w, 0.38, PURPLE)
    txt(s, x+0.18, y+0.075, w-0.3, 0.26, [(label, 13.5, True, WHITE)])

def bullet(s, x, y, w, text, size=12.5, bold=True, color=PURPLE, dot=GREEN, h=0.5):
    rect(s, x, y+0.075, 0.09, 0.09, dot)
    txt(s, x+0.26, y, w, h, [(text, size, bold, color, 0, 1.22)])


def chip(s, x, y, w, h, text, fill, fg=WHITE, size=9.5):
    sh = rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.5)
    txt(s, x, y+(h-0.17)/2, w, 0.20, [(text, size, True, fg)], align=PP_ALIGN.CENTER)
    return sh

def kicker_header(s, kicker, title, sub=None):
    rect(s, 0, 0, 13.333, 0.09, PURPLE)
    rect(s, 0, 0.09, 1.60, 0.04, GREEN)
    s.shapes.add_picture(LOGO, Inches(11.55), Inches(0.30), Inches(1.42), Inches(0.30))
    txt(s, 0.55, 0.26, 10.0, 0.28, [(kicker, 10.5, True, GREEN)])
    txt(s, 0.55, 0.55, 10.6, 0.50, [(title, 25, True, PURPLE)])
    if sub:
        txt(s, 0.55, 1.08, 11.2, 0.32, [(sub, 12.5, False, GRAY)])
        return 1.52
    return 1.20
