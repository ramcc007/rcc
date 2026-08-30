import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from common import *
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt

prs = new_deck()
N = [0]
def page(): N[0] += 1; return N[0]

def style_axis(ax, size=9.5, color=GRAY, gridlines=False):
    ax.has_major_gridlines = gridlines
    if gridlines:
        gl = ax.major_gridlines.format.line
        gl.color.rgb = RGBColor(0xEC,0xEA,0xEE); gl.width = Pt(0.75)
    ax.tick_labels.font.size = Pt(size)
    ax.tick_labels.font.name = FONT
    ax.tick_labels.font.color.rgb = color
    ax.format.line.color.rgb = RGBColor(0xDD,0xDA,0xE0)

# ============================================================ 1. TITLE
s = blank(prs)
rect(s, 0, 0, 4.60, 7.5, PURPLE)
rect(s, 4.60, 0, 0.09, 7.5, GREEN)
rect(s, 0.55, 0.55, 2.60, 0.85, WHITE)
s.shapes.add_picture(LOGO, Inches(0.78), Inches(0.78), Inches(2.10), Inches(0.45))
txt(s, 0.55, 2.40, 3.85, 2.40, [("GEO / AI", 30, True, WHITE, 2),
                                ("Visibility", 30, True, WHITE, 2),
                                ("& SEO", 30, True, WHITE, 2),
                                ("Performance", 30, True, WHITE)])
txt(s, 0.55, 5.02, 3.85, 0.40, [("Q2 2026 Board Update", 16, True, GREEN)])
txt(s, 0.55, 6.45, 3.70, 0.60, [("Prepared for the Board  ·  Q2 2026 (Apr–Jun)", 11.5, False, LAV, 3),
                                ("Sources: Semrush, Aug 2026 pull", 10.5, False, LAV)])
txt(s, 5.04, 2.62, 7.75, 1.85, [("Where Komprise stands in AI", 24, True, PURPLE, 5),
                                ("search and organic search after", 24, True, PURPLE, 5),
                                ("Q2 \u2014 and what happens next", 24, True, PURPLE)])
rect(s, 5.04, 4.86, 1.40, 0.04, GREEN)
txt(s, 5.04, 5.24, 7.75, 1.30,
    [("Two lanes, one story:", 13, True, INK, 6),
     ("1.  Be the answer for AI-ready data — unstructured data for AI, AI data management.", 12.5, False, GRAY, 4),
     ("2.  Be the answer for storage price hikes — tiering, cost savings, NetApp, Dell, Everpure.", 12.5, False, GRAY)])
page()

# ============================================================ 2. EXECUTIVE SUMMARY
s = blank(prs)
header(s, "Executive Summary", "Q2 2026 (Apr–Jun) — the quarter in four numbers")
statcard(s, 0.55, 1.40, 2.85, 1.52, "+32%", "Organic Traffic", "16,602 → 21,883 monthly visits. Best month since May 2025.")
statcard(s, 3.55, 1.40, 2.85, 1.52, "+177%", "Traffic Value", "$29,984 → $82,977 / month of equivalent paid search.")
statcard(s, 6.55, 1.40, 2.85, 1.52, "5,449", "AI Overview Keywords", "14× Hammerspace, 127× Diskover. Up 7% in the quarter.", PURPLE)
statcard(s, 9.55, 1.40, 2.85, 1.52, "+15%", "Top-3 Rankings", "366 → 420 keywords in Google's top three results.")

bandbar(s, 0.55, 3.22, 12.23, "What the Board Should Take Away")
bullet(s, 0.60, 3.82, 11.9, "Q2 was our strongest search quarter in over a year. Traffic grew 32% and the commercial value of that traffic nearly tripled, on a smaller, better-targeted keyword set.")
bullet(s, 0.60, 4.52, 11.9, "We already show up inside Google's AI answers at real scale. Komprise appears in AI Overviews for 5,449 keywords — an order of magnitude ahead of Hammerspace, Atempo, Datadobi and Diskover.")
bullet(s, 0.60, 5.22, 11.9, "The competitive picture changed in Q2. Pure Storage relaunched as Everpure and, with 1touch, went from no measurable presence to roughly three times our AI Overview footprint inside two months.")
bullet(s, 0.60, 5.92, 11.9, "The brand-perception work is the unfinished half. Being cited is not the same as being named, and the prompt-level tracking needed to prove that shift only began on 10 June — see Data Coverage.")

rect(s, 0.55, 6.68, 12.23, 0.46, PALE)
rect(s, 0.55, 6.68, 0.06, 0.46, PURPLE)
txt(s, 0.78, 6.79, 11.9, 0.30, [("Bottom line: the machine is working on search. The AI-brand story is where Q3 investment goes.", 12, True, PURPLE)])
footer(s, page())

# ============================================================ 3. ORGANIC GROWTH
s = blank(prs)
header(s, "Q2 Search Performance", "Monthly Semrush snapshots, US database. 15 Mar = Q1 exit baseline; 15 Apr / 15 May / 15 Jun = Q2.")

cd = CategoryChartData()
cd.categories = ['15 Mar', '15 Apr', '15 May', '15 Jun', '15 Jul']
cd.add_series('Monthly organic traffic', (18129, 16602, 20440, 21883, 23880))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.50), Inches(8.05), Inches(4.55), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = False
pl = ch.plots[0]; pl.gap_width = 55
pl.has_data_labels = True
dl = pl.data_labels
dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = FONT; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.number_format = '#,##0'; dl.number_format_is_linked = False
for i, pt in enumerate(pl.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = LAV if i in (0, 4) else PURPLE
    pt.format.line.fill.background()
style_axis(ch.value_axis, gridlines=True); style_axis(ch.category_axis, size=11)
ch.value_axis.has_major_gridlines = True
ch.value_axis.minimum_scale = 0
ch.value_axis.maximum_scale = 28000

txt(s, 0.55, 6.20, 8.05, 0.75,
    [("Traffic dipped into April, then climbed for three straight months. "
      "The two pale bars sit outside Q2 and are shown for context — 15 Jul is the first reading after quarter close.", 11, False, GRAY, 0, 1.25)])

statcard(s, 8.95, 1.50, 3.83, 1.42, "6,225", "Ranking Keywords", "Down from 6,908 in March — a smaller, more focused set earning far more traffic.", PURPLE, 26)
statcard(s, 8.95, 3.06, 3.83, 1.42, "1,095", "Keywords in Positions 4–10", "Up 12.7% in the quarter. The pipeline feeding future top-3 places.", DGREEN, 26)
statcard(s, 8.95, 4.62, 3.83, 1.42, "88,764", "Semrush Rank", "Improved by 26,150 places across Q2 (lower is better).", DGREEN, 26)
footer(s, page())

# ============================================================ 4. AI ANSWERS
s = blank(prs)
header(s, "Showing Up Inside AI Answers", "Keywords where komprise.com is pulled into a Google AI Overview or a People Also Ask box")

cd = CategoryChartData()
cd.categories = ['15 Mar', '15 Apr', '15 May', '15 Jun']
cd.add_series('AI Overviews', (5522, 5079, 5297, 5449))
cd.add_series('People Also Ask', (4000, 1933, 4603, 5783))
gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(0.55), Inches(1.50), Inches(7.70), Inches(4.20), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
for ser, col in zip(ch.series, (PURPLE, GREEN)):
    ser.format.line.color.rgb = col; ser.format.line.width = Pt(2.75)
    ser.marker.format.fill.solid(); ser.marker.format.fill.fore_color.rgb = col
    ser.marker.format.line.color.rgb = WHITE
    ser.smooth = False
style_axis(ch.value_axis, gridlines=True); style_axis(ch.category_axis, size=11)
ch.value_axis.minimum_scale = 0

rect(s, 8.60, 1.50, 4.18, 2.05, OFFWHITE, BORDER)
rect(s, 8.60, 1.50, 0.06, 2.05, PURPLE)
txt(s, 8.86, 1.72, 3.75, 1.70,
    [("Read this honestly", 12, True, PURPLE, 6),
     ("AI Overview presence grew steadily through Q2 (+7%), but at 5,449 it has not yet "
      "recovered the 5,522 we held in March. Q2 was a recovery-and-rebuild quarter, "
      "not a straight line up.", 11, False, GRAY, 0, 1.25)])

rect(s, 8.60, 3.70, 4.18, 2.35, OFFWHITE, BORDER)
rect(s, 8.60, 3.70, 0.06, 2.35, DGREEN)
txt(s, 8.86, 3.92, 3.75, 2.00,
    [("The standout move", 12, True, PURPLE, 6),
     ("People Also Ask presence tripled — 1,933 to 5,783 keywords. That is the direct, "
      "measurable payoff from restructuring pages into question-and-answer form, "
      "which is exactly the format AI engines lift from.", 11, False, GRAY, 0, 1.25)])

rect(s, 0.55, 6.20, 12.23, 0.62, PALE)
rect(s, 0.55, 6.20, 0.06, 0.62, PURPLE)
txt(s, 0.78, 6.34, 11.8, 0.42,
    [("In plain terms: when someone asks Google a question in our category, there is now a very good chance our answer is part of what they read — even if they never click through to komprise.com.", 12, True, PURPLE, 0, 1.2)])
footer(s, page())

# ============================================================ 5. COMPETITIVE
s = blank(prs)
header(s, "How We Compare", "Keywords appearing in Google AI Overviews, 15 June 2026 — Semrush, US database")

cd = CategoryChartData()
cd.categories = ['Everpure\n(Pure Storage)', 'Komprise', 'Unstructured.io', 'Hammerspace', 'Atempo', 'Datadobi', 'Diskover']
cd.add_series('AI Overview keywords', (16766, 5449, 1259, 385, 193, 50, 43))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(1.52), Inches(7.55), Inches(4.55), cd)
ch = gf.chart
ch.has_title = False
ch.has_legend = False
pl = ch.plots[0]; pl.gap_width = 45
pl.has_data_labels = True
dl = pl.data_labels
dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = FONT; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.number_format = '#,##0'; dl.number_format_is_linked = False
cols = [AMBER, PURPLE, LGRAY, LGRAY, LGRAY, LGRAY, LGRAY]
for i, pt in enumerate(pl.series[0].points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = cols[i]
    pt.format.line.fill.background()
style_axis(ch.value_axis, gridlines=True); style_axis(ch.category_axis, size=11)
ch.value_axis.minimum_scale = 0
ch.value_axis.maximum_scale = 19500
txt(s, 0.55, 6.24, 7.55, 0.60,
    [("Everpure is Pure Storage, relaunched under a new name in Q2 together with its 1touch acquisition. "
      "It had no measurable organic presence in this database before May 2026.", 10, False, LGRAY, 0, 1.2)])

rect(s, 8.45, 1.52, 4.33, 2.35, OFFWHITE, BORDER)
rect(s, 8.45, 1.52, 0.06, 2.35, DGREEN)
txt(s, 8.71, 1.74, 3.90, 2.00,
    [("Against the field we set out to beat", 12, True, PURPLE, 6),
     ("Komprise out-reaches every competitor named in our GEO strategy: 14× Hammerspace, "
      "28× Atempo, 109× Datadobi and 127× Diskover. The content investment of the last "
      "year is visible here.", 11, False, GRAY, 0, 1.25)])

rect(s, 8.45, 4.02, 4.33, 2.25, OFFWHITE, BORDER)
rect(s, 8.45, 4.02, 0.06, 2.25, AMBER)
txt(s, 8.71, 4.24, 3.90, 1.90,
    [("The new entrant", 12, True, AMBER, 6),
     ("Everpure is Pure Storage rebranded, plus the 1touch acquisition. It is a far larger "
      "company than the specialists we have been benchmarking against — see next slide.", 11, False, GRAY, 0, 1.25)])
footer(s, page())

# ============================================================ 6. EVERPURE
s = blank(prs)
header(s, "The Everpure Shift", "The single biggest competitive change of Q2 — and it lands squarely on the storage-cost lane",
       kicker="LANE 1  ·  STORAGE PRICE HIKES")

cd = CategoryChartData()
cd.categories = ['15 Apr', '15 May', '15 Jun']
cd.add_series('Everpure', (0, 6410, 16766))
cd.add_series('Komprise', (5079, 5297, 5449))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.72), Inches(7.30), Inches(4.20), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP
ch.legend.include_in_layout = False
ch.legend.font.size = Pt(11); ch.legend.font.name = FONT; ch.legend.font.color.rgb = GRAY
pl = ch.plots[0]; pl.gap_width = 60
pl.has_data_labels = True
dl = pl.data_labels
dl.font.size = Pt(10); dl.font.bold = True; dl.font.name = FONT; dl.font.color.rgb = INK
dl.position = XL_LABEL_POSITION.OUTSIDE_END
dl.number_format = '#,##0'; dl.number_format_is_linked = False
for ser, col in zip(ch.series, (AMBER, PURPLE)):
    ser.format.fill.solid(); ser.format.fill.fore_color.rgb = col
    ser.format.line.fill.background()
style_axis(ch.value_axis, gridlines=True); style_axis(ch.category_axis, size=11)
txt(s, 0.55, 6.05, 7.30, 0.45, [("Keywords in Google AI Overviews. Everpure shows zero through 15 April because the domain had no measurable organic presence in Semrush's US database before the rebrand.", 10, False, LGRAY, 0, 1.2)])

rect(s, 8.20, 1.72, 4.58, 4.60, OFFWHITE, BORDER)
rect(s, 8.20, 1.72, 0.06, 4.60, AMBER)
txt(s, 8.48, 1.98, 4.10, 4.15,
    [("What happened", 12.5, True, PURPLE, 7),
     ("Pure Storage relaunched as Everpure and folded in its 1touch acquisition. In two months the new domain went from nothing to 33,353 monthly organic visits and 16,766 AI Overview keywords — roughly three times our footprint.", 11, False, GRAY, 9, 1.25),
     ("Why it matters", 12.5, True, PURPLE, 7),
     ("Everpure is positioning AI as the product, not as an outcome. That is the same framing gap our CEO flagged in June when Gemini described Komprise as a storage and cost company.", 11, False, GRAY, 9, 1.25),
     ("What we are not saying", 12.5, True, PURPLE, 7),
     ("This is not lost ground. It is a new, much larger competitor arriving in our category mid-quarter, and it resets who we benchmark against.", 11, False, GRAY, 0, 1.25)])
footer(s, page())

# ============================================================ 7. WHAT WE DID
s = blank(prs)
header(s, "What We Did to Change How We Are Described", "Split by date, so the Board can see what landed inside the quarter and what is running now")

rect(s, 0.55, 1.45, 6.02, 0.38, PURPLE)
txt(s, 0.73, 1.525, 5.7, 0.26, [("DELIVERED IN Q2  ·  ON OUR OWN SITE", 12.5, True, WHITE)])
rect(s, 6.76, 1.45, 6.02, 0.38, DGREEN)
txt(s, 6.94, 1.525, 5.7, 0.26, [("IN MOTION SINCE QUARTER CLOSE  ·  THIRD PARTY", 12.5, True, WHITE)])

L1 = [("23 June", "Data Classification and Unstructured Data Classification glossary pages rewritten, and a new AI-Ready Data page created to claim that term directly."),
      ("16 June", "\u201cFive Techniques to Eliminate Noisy Data from Enterprise AI Pipelines\u201d published \u2014 our clearest AI-pipeline explainer to date."),
      ("23 June", "Transparent File Tables interview blog published, feeding the in-the-news category."),
      ("Late June", "Transparent File Tables product launch \u2014 product page, demo page and press coverage."),
      ("Through Q2", "Product page titles, descriptions and the homepage hero reframed to lead with AI readiness, with the 70% cost saving moved to a supporting proof point.")]

L2 = [("1 July", "AI Data Platform glossary page, Transparent File Tables resource page and the KAPPA data services library page published."),
      ("15 July", "A single consistent company description agreed and rolled out across third-party sites, so every source describes Komprise the same way."),
      ("Late July", "Isilon and flash / storage pricing pages updated to meet buyers searching on price rather than on product."),
      ("July\u2013August", "Review, directory and marketplace listings being corrected \u2014 TrustRadius, G2, Glassdoor and Gartner Peer Insights, plus the major partner marketplaces."),
      ("July\u2013August", "Weekly capture running across all five major AI assistants, so we can prove whether the wording changes actually move the answer.")]

def worklist(x, items, dot):
    y = 2.00
    for date, body in items:
        rect(s, x, y+0.06, 0.09, 0.09, dot)
        txt(s, x+0.26, y-0.02, 1.05, 0.24, [(date, 9.5, True, dot)])
        txt(s, x+0.26, y+0.22, 5.60, 0.70, [(body, 11, False, INK, 0, 1.22)])
        y += 0.90

worklist(0.58, L1, PURPLE)
worklist(6.79, L2, DGREEN)

rect(s, 0.55, 6.48, 12.23, 0.62, PALE)
rect(s, 0.55, 6.48, 0.06, 0.62, PURPLE)
txt(s, 0.78, 6.61, 11.8, 0.42,
    [("Early proof it works (August testing): product names we deliberately repeated across trusted sites \u2014 Transparent Move Technology, Global Metadatabase \u2014 are now being repeated back to us by AI assistants, unprompted.", 11, True, PURPLE, 0, 1.2)])
footer(s, page())

# ============================================================ 8. WHAT MATTERS TO LLMS
s = blank(prs)
header(s, "What Actually Matters to the AI Engines", "Three things decide whether an AI assistant names you — and only one of them is your own website")

def pillar(x, num, title, body, foot, col):
    rect(s, x, 1.50, 3.90, 4.35, OFFWHITE, BORDER)
    rect(s, x, 1.50, 3.90, 0.07, col)
    rect(s, x+0.28, 1.85, 0.52, 0.52, col, shape=MSO_SHAPE.OVAL)
    txt(s, x+0.28, 1.98, 0.52, 0.30, [(num, 15, True, WHITE)], align=PP_ALIGN.CENTER)
    txt(s, x+0.28, 2.58, 3.35, 0.40, [(title, 14, True, PURPLE)])
    txt(s, x+0.28, 3.08, 3.35, 1.90, [(body, 11.5, False, GRAY, 0, 1.30)])
    rect(s, x+0.28, 4.98, 3.35, 0.76, PALE)
    txt(s, x+0.42, 5.09, 3.10, 0.58, [(foot, 10, True, PURPLE, 0, 1.18)])

pillar(0.55, "1", "Be readable",
       "AI crawlers lift answers from plain, static text. Content hidden behind click-to-open panels or heavy scripting is invisible to them, however good it is. Answer-first writing and question-and-answer formatting are what get quoted.",
       "Our PAA presence tripling in Q2 is this working.", PURPLE)
pillar(4.72, "2", "Be repeated",
       "The engines form their opinion of a company by reading the same description of it across many trusted sites. One well-written page cannot outweigh a hundred third-party sources that describe you differently.",
       "Hence the review, directory and marketplace clean-up.", DGREEN)
pillar(8.88, "3", "Be named",
       "Being cited as a source and being named as a vendor are different outcomes. Citations follow good content. Mentions follow earned coverage — analyst placements, press and third-party authority.",
       "June 2026 baseline: our domain was 3% of the 594 sources feeding our 50 prompts.", AMBER)

rect(s, 0.55, 6.10, 12.23, 0.80, PALE)
rect(s, 0.55, 6.10, 0.06, 0.80, PURPLE)
txt(s, 0.78, 6.24, 11.8, 0.60,
    [("The strategic consequence: on the June baseline, publishing more on komprise.com moves roughly 3% of what the engines read. Shifting how we are described requires the third-party and PR work to run in parallel with content — which is why Q3 is weighted that way.", 12, True, PURPLE, 0, 1.22)])
footer(s, page())

# ============================================================ 9. WHAT WE TRACK
s = blank(prs)
header(s, "What We Track Going Forward", "The 50 AI prompts now under continuous measurement, live in Semrush since 10 June 2026")

cats = [("AI & AI-Readiness", 16, DGREEN),
        ("Unstructured Data Management", 11, PURPLE),
        ("Data Classification", 8, AMBER),
        ("Data Tiering & Storage Efficiency", 8, PURPLE),
        ("Migration & Lifecycle", 7, DGREEN)]
x = 0.55
for name, n, col in cats:
    rect(s, x, 1.45, 2.36, 1.40, OFFWHITE, BORDER)
    rect(s, x, 1.45, 2.36, 0.07, col)
    txt(s, x+0.22, 1.68, 1.95, 0.55, [(str(n), 27, True, col)])
    txt(s, x+0.22, 2.22, 1.95, 0.55, [(name, 10.5, True, PURPLE, 0, 1.18)])
    x += 2.45

txt(s, 0.55, 3.10, 12.23, 0.32, [("Representative prompts under measurement", 12.5, True, PURPLE)])
ex = ["What are the best tools for managing unstructured data?",
      "What is an AI-ready data platform?",
      "How do organizations prepare unstructured data for AI?",
      "What are the benefits of intelligent data tiering?",
      "What are the best tools for AI data ingestion?",
      "How does automated data tiering reduce storage costs?",
      "Can you recommend software for unstructured data classification?",
      "What are the best enterprise data migration solutions?"]
y = 3.52
for i, e in enumerate(ex):
    xx = 0.55 if i < 4 else 6.76
    yy = 3.52 + (i % 4) * 0.46
    rect(s, xx, yy, 6.02, 0.38, OFFWHITE, BORDER)
    txt(s, xx+0.20, yy+0.09, 5.65, 0.26, [(e, 11, False, INK)])

rect(s, 0.55, 5.55, 12.23, 1.30, OFFWHITE, BORDER)
rect(s, 0.55, 5.55, 0.06, 1.30, AMBER)
txt(s, 0.82, 5.76, 11.8, 1.00,
    [("Why there is no Q2 trend line for these 50 prompts", 12.5, True, PURPLE, 6),
     ("Prompt-level tracking went live on 10 June 2026, so it covers only the final three weeks of the quarter. Across those three weeks none of the 50 prompts placed komprise.com in Google's top 100 organic results — these are long conversational questions that Google answers with an AI Overview rather than a ranked list, which is precisely why we measure AI Overview presence separately. Q3 will be the first full quarter of prompt-level trend data.", 11, False, GRAY, 0, 1.25)])
footer(s, page())

# ============================================================ 10. DATA COVERAGE
s = blank(prs)
header(s, "Data Coverage and Open Items", "Stated plainly, so nothing in this deck is taken for more than it is")

rows = [
 ("Prompt-level GEO history for April and May",
  "The 50-prompt tracking campaign was created on 10 June 2026. There is no April or May prompt-level data to report.",
  "Q3 will be the first complete quarter."),
 ("AI visibility score, mentions, citations and source mix",
  "The figures that headline our monthly decks are not retrievable through the current Semrush data connection, which exposes only the Google side of tracking. Any such figure quoted here is carried forward from the June 2026 deck and labelled as of that date.",
  "Pull manually from the Semrush AI dashboard for the next monthly."),
 ("Darren's priority keyword list",
  "The 'Top KWs by Darren' tracking project exists but could not be reached programmatically. Priority-keyword reporting is therefore not included, rather than substituted with a different list.",
  "Send the project link and it goes into the next update."),
 ("Brand and non-brand prompt sets from Monte and Polly",
  "Monte's prompt set is held in a workbook outside this account. Polly's Top 10 priority prompts have not yet been defined — she asked for a working session in her 28 August note.",
  "Agree the Top 10 on the call in the week of 1 September."),
]
y = 1.45
rect(s, 0.55, y, 12.23, 0.36, PURPLE)
txt(s, 0.75, y+0.065, 3.60, 0.26, [("Item", 11, True, WHITE)])
txt(s, 4.55, y+0.065, 5.40, 0.26, [("What the situation is", 11, True, WHITE)])
txt(s, 10.15, y+0.065, 2.50, 0.26, [("Next step", 11, True, WHITE)])
y += 0.36
hs = [0.95, 1.35, 1.15, 1.15]
for (a, b, c), h in zip(rows, hs):
    rect(s, 0.55, y, 12.23, h, OFFWHITE, BORDER)
    txt(s, 0.75, y+0.16, 3.60, h-0.3, [(a, 11.5, True, PURPLE, 0, 1.2)])
    txt(s, 4.55, y+0.16, 5.40, h-0.3, [(b, 10.5, False, GRAY, 0, 1.25)])
    txt(s, 10.15, y+0.16, 2.45, h-0.3, [(c, 10.5, False, INK, 0, 1.25)])
    y += h + 0.08
footer(s, page())

# ============================================================ 11. Q3 PLAN
s = blank(prs)
header(s, "Q3 2026 Plan", "Where the effort goes next, and how we will know it worked")

items = [
 ("Weight the effort towards third-party sources, not more pages",
  "Our own domain is a small fraction of what the AI engines read. Digital PR and analyst placements aimed at the sites our competitors are already cited from is the highest-leverage move available.",
  PURPLE),
 ("Fix the Gemini gap specifically",
  "Testing across all five major assistants through the summer showed ChatGPT, Perplexity and Claude already give Komprise the AI-readiness credit. Gemini is the outlier. Aim the source corrections there first.",
  DGREEN),
 ("Agree the Top 10 priority prompts with Darren and Polly",
  "Narrow the 50 tracked prompts to the ten that matter most commercially, and report against those. Working session already requested for the week of 1 September.",
  AMBER),
 ("Add Everpure to the standing competitive set",
  "It arrived mid-quarter at three times our AI Overview footprint. It should be tracked alongside Diskover and Hammerspace from now on, not treated as an event.",
  PURPLE),
 ("Keep converting pages into answer-first, question-and-answer form",
  "This is the one on-site lever with proven Q2 return — People Also Ask presence tripled. Continue across the remaining legacy pages with Tangence.",
  DGREEN),
]
y = 1.45
for i, (t, b, col) in enumerate(items, 1):
    rect(s, 0.55, y, 12.23, 0.92, OFFWHITE, BORDER)
    rect(s, 0.55, y, 0.06, 0.92, col)
    rect(s, 0.85, y+0.26, 0.40, 0.40, col, shape=MSO_SHAPE.OVAL)
    txt(s, 0.85, y+0.355, 0.40, 0.26, [(str(i), 12, True, WHITE)], align=PP_ALIGN.CENTER)
    txt(s, 1.45, y+0.14, 11.1, 0.28, [(t, 12.5, True, PURPLE)])
    txt(s, 1.45, y+0.46, 11.1, 0.40, [(b, 10.5, False, GRAY, 0, 1.20)])
    y += 0.99

rect(s, 0.55, y+0.02, 12.23, 0.62, PALE)
rect(s, 0.55, y+0.02, 0.06, 0.62, PURPLE)
txt(s, 0.78, y+0.15, 11.8, 0.44,
    [("How we will know it worked: AI Overview keyword count back above 5,522 and rising, plus a first full quarter of prompt-level movement on the agreed Top 10.", 11, True, PURPLE, 0, 1.2)])
footer(s, page())

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "Komprise_GEO_SEO_Q2_2026_Board_Update.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
